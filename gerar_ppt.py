import csv, json, io, os
from collections import defaultdict
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colours ──────────────────────────────────────────────────────
BG     = RGBColor(0x0f, 0x11, 0x17)
CARD   = RGBColor(0x1a, 0x1d, 0x27)
ACCENT = RGBColor(0x63, 0x66, 0xf1)
WHITE  = RGBColor(0xff, 0xff, 0xff)
MUTED  = RGBColor(0x94, 0xa3, 0xb8)
GREEN  = RGBColor(0x22, 0xc5, 0x5e)
RED    = RGBColor(0xef, 0x44, 0x44)
YELLOW = RGBColor(0xf5, 0x9e, 0x0b)
ROW_ALT = RGBColor(0x1e, 0x21, 0x30)

PALETTE_HEX = [
    '#6366f1','#22c55e','#f59e0b','#ef4444','#06b6d4','#ec4899',
    '#a855f7','#14b8a6','#f97316','#84cc16','#3b82f6','#e11d48',
    '#10b981','#f472b6','#fb923c','#a3e635','#38bdf8','#c084fc',
    '#fde68a','#6ee7b7'
]

def hex_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) / 255 for i in (0, 2, 4))

PALETTE_RGB = [hex_rgb(c) for c in PALETTE_HEX]

PERIODS = ['Pre_01-16Mar','S1_17-23Mar','S2_24-30Mar','S3_31Mar-06Abr','S4_07-13Abr','S5_14-16Abr']
PLABELS = ['1-16/Mar', 'S1 17-23/Mar', 'S2 24-30/Mar', 'S3 31Mar-6/Abr', 'S4 7-13/Abr', 'S5 14-16/Abr']
PLABELS_SHORT = ['1-16/Mar', 'S1', 'S2', 'S3', 'S4', 'S5']

# ── load data ────────────────────────────────────────────────────
rows = []
with open("C:/claudinho/outgoing_evolution_data.csv", encoding='utf-8') as f:
    for r in csv.DictReader(f):
        rows.append({
            'cust': r['CUS_CUST_ID'], 'semana': r['semana'],
            'processo': r['processo'], 'cdu': r['cdu'],
            'solucao': r['solucao'], 'qtd': int(r['qtd'])
        })

data   = defaultdict(lambda: defaultdict(lambda: defaultdict(lambda: defaultdict(lambda: defaultdict(int)))))
totals = defaultdict(lambda: defaultdict(int))
for r in rows:
    data[r['cust']][r['semana']][r['processo']][r['cdu']][r['solucao']] += r['qtd']
    totals[r['cust']][r['semana']] += r['qtd']

customers = sorted(data.keys(), key=lambda x: -sum(totals[x].values()))
wt = {c: [totals[c].get(p, 0) for p in PERIODS] for c in customers}
grand_total = sum(sum(wt[c]) for c in customers)

def proc_weekly(c):
    pt = defaultdict(lambda: [0]*len(PERIODS))
    for pi, p in enumerate(PERIODS):
        for proc, cdus in data[c][p].items():
            pt[proc][pi] += sum(sum(v.values()) for v in cdus.values())
    return dict(sorted(pt.items(), key=lambda x: -sum(x[1])))

# ── matplotlib style ─────────────────────────────────────────────
plt.rcParams.update({
    'figure.facecolor': '#0f1117', 'axes.facecolor': '#1a1d27',
    'axes.edgecolor': '#2a2d3a', 'axes.labelcolor': '#94a3b8',
    'xtick.color': '#64748b', 'ytick.color': '#64748b',
    'text.color': '#e2e8f0', 'grid.color': '#1e2130', 'grid.linewidth': 0.5,
    'font.family': 'DejaVu Sans', 'font.size': 9,
})

def fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

# ── PPT helpers ───────────────────────────────────────────────────
def set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, text, l, t, w, h, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def add_rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def kpi_box(slide, l, t, w, h, label, value, delta=None, delta_up=None):
    add_rect(slide, l, t, w, h, CARD)
    add_text(slide, label, l+0.1, t+0.07, w-0.2, 0.22, size=7, color=MUTED)
    add_text(slide, str(value), l+0.1, t+0.27, w-0.2, 0.38, size=15, bold=True, color=WHITE)
    if delta:
        dcol = RED if delta_up else GREEN
        add_text(slide, delta, l+0.1, t+0.63, w-0.2, 0.22, size=8, color=dcol)

def draw_table(slide, headers, rows_data, col_widths, start_x, start_y, row_h=0.22, hdr_color=ACCENT):
    col_x = [start_x]
    for w in col_widths[:-1]:
        col_x.append(col_x[-1] + w)
    total_w = sum(col_widths)
    # header
    add_rect(slide, start_x, start_y, total_w, row_h, hdr_color)
    for j, (h, x) in enumerate(zip(headers, col_x)):
        add_text(slide, h, x+0.04, start_y+0.03, col_widths[j]-0.05, row_h-0.04, size=8, bold=True, color=WHITE)
    # rows
    for i, row in enumerate(rows_data):
        y = start_y + row_h * (i + 1)
        if y > 7.25:
            break
        bg = ROW_ALT if i % 2 == 0 else CARD
        add_rect(slide, start_x, y, total_w, row_h, bg)
        for j, (val, x) in enumerate(zip(row, col_x)):
            col = row[-1] if isinstance(row[-1], RGBColor) else WHITE
            # last element can be a colour hint — skip it
            if j < len(headers):
                add_text(slide, str(val)[:80], x+0.04, y+0.03, col_widths[j]-0.05, row_h-0.04, size=8, color=WHITE)

# ════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── SLIDE 1: CAPA ────────────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_rect(sl, 0, 0, 0.08, 7.5, ACCENT)
add_text(sl, 'OUTGOING',         0.5, 1.6, 12, 1.1, size=52, bold=True,  color=ACCENT)
add_text(sl, 'Analise de Contatos Respondidos por CUS_CUST_ID',
         0.5, 2.85, 12, 0.6, size=20, color=WHITE)
add_text(sl,
    'Equipes: BR_ME_Pre-Despacho_Offline  |  BR_MLPostVenda_Offline  |  BR_MLPreVenda_Offline  |  BR_Prustomer_Offline\n'
    'BR_ME_Sellers_Longtail  |  BR_Publicaciones_Sellers_Longtail  |  BR_Ventas_Sellers_Longtail',
    0.5, 3.65, 12, 0.8, size=11, color=MUTED)
add_text(sl, 'ci_event_name = OUTGOING_CONTACT  |  OUTGOING_FIRST_CONTACT',
         0.5, 4.5, 12, 0.35, size=11, color=MUTED)
add_text(sl, 'Periodo: 01/Mar - 16/Abr/2026   |   Top 20 Clientes',
         0.5, 5.0, 12, 0.4, size=14, bold=True, color=YELLOW)
add_text(sl, f'Total Geral: {grand_total:,} outgoing',
         0.5, 5.55, 8, 0.4, size=15, bold=True, color=WHITE)
add_text(sl, 'Gerado em 16/04/2026', 0.5, 6.9, 6, 0.3, size=9, color=MUTED)

# ── SLIDE 2: VISÃO GERAL ─────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_text(sl, 'Visao Geral - Evolucao Semanal Top 20', 0.3, 0.1, 12, 0.42, size=16, bold=True)
add_text(sl, '01/Mar - 16/Abr/2026  |  Stacked por cliente', 0.3, 0.5, 12, 0.28, size=10, color=MUTED)

# Stacked bar all customers
fig, ax = plt.subplots(figsize=(8.2, 3.2))
bottoms = np.zeros(len(PERIODS))
for i, c in enumerate(customers):
    vals = np.array(wt[c])
    ax.bar(PLABELS_SHORT, vals, bottom=bottoms, color=PALETTE_RGB[i % len(PALETTE_RGB)], label=c, width=0.6)
    bottoms += vals
ax.set_ylabel('Qtd Outgoing')
ax.yaxis.grid(True)
ax.set_axisbelow(True)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
buf = fig_to_stream(fig)
sl.shapes.add_picture(buf, Inches(0.3), Inches(0.88), Inches(8.0), Inches(3.4))

# Ranking top 10 horizontal
fig2, ax2 = plt.subplots(figsize=(4.5, 3.2))
top10 = customers[:10]
t10_vals = [sum(wt[c]) for c in top10][::-1]
t10_labs = top10[::-1]
t10_cols = [PALETTE_RGB[i % len(PALETTE_RGB)] for i in range(len(top10))][::-1]
ax2.barh(t10_labs, t10_vals, color=t10_cols, height=0.6)
ax2.xaxis.grid(True)
ax2.set_axisbelow(True)
ax2.spines['top'].set_visible(False)
ax2.spines['right'].set_visible(False)
ax2.set_xlabel('Total Outgoing')
for i, v in enumerate(t10_vals):
    ax2.text(v + max(t10_vals)*0.01, i, str(v), va='center', fontsize=8, color='#e2e8f0')
buf2 = fig_to_stream(fig2)
sl.shapes.add_picture(buf2, Inches(8.5), Inches(0.88), Inches(4.5), Inches(3.4))

# KPIs
total_mar = sum(sum(wt[c][:3]) for c in customers)
total_abr = sum(sum(wt[c][3:]) for c in customers)
total_s4  = sum(wt[c][4] for c in customers)
total_s5  = sum(wt[c][5] for c in customers)
kpi_items = [
    ('Total Geral',          f'{grand_total:,}'),
    ('Marco/26',             f'{total_mar:,}'),
    ('Abril/26 (ate 16)',    f'{total_abr:,}'),
    ('S4 (7-13/Abr)',        f'{total_s4:,}'),
    ('S5 (14-16/Abr)',       f'{total_s5:,}'),
    ('Top Cliente',          customers[0]),
]
for idx, (lbl, val) in enumerate(kpi_items):
    kpi_box(sl, 0.3 + idx * 2.12, 4.5, 2.0, 1.05, lbl, val)

# ── SLIDE 3: MES VS MES ───────────────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_text(sl, 'Mes vs Mes - Top 20 Clientes', 0.3, 0.1, 12, 0.42, size=16, bold=True)
add_text(sl, 'Marco/2026 vs Abril/2026 (ate 16)', 0.3, 0.5, 12, 0.28, size=10, color=MUTED)

fig, ax = plt.subplots(figsize=(12.5, 3.9))
x = np.arange(len(customers))
bar_w = 0.35
mar_vals = [sum(wt[c][:3]) for c in customers]
abr_vals = [sum(wt[c][3:]) for c in customers]
ax.bar(x - bar_w/2, mar_vals, bar_w, label='Mar/26', color='#6366f1', alpha=0.85)
ax.bar(x + bar_w/2, abr_vals, bar_w, label='Abr/26', color='#22c55e', alpha=0.85)
ax.set_xticks(x)
ax.set_xticklabels(customers, rotation=35, ha='right', fontsize=7.5)
ax.legend(fontsize=9)
ax.yaxis.grid(True)
ax.set_axisbelow(True)
ax.spines['top'].set_visible(False)
ax.spines['right'].set_visible(False)
buf = fig_to_stream(fig)
sl.shapes.add_picture(buf, Inches(0.3), Inches(0.88), Inches(12.7), Inches(4.0))

# Table
col_w = [0.35, 1.55, 0.9, 1.0, 0.9, 0.75]
col_x = [0.3]
for w in col_w[:-1]:
    col_x.append(col_x[-1] + w)
tw = sum(col_w)
row_h = 0.215
sy = 5.08
hdrs = ['#', 'CUS_CUST_ID', 'Mar/26', 'Abr/26', 'Total', 'Var%']
add_rect(sl, 0.3, sy, tw, row_h, ACCENT)
for j, (h, x2) in enumerate(zip(hdrs, col_x)):
    add_text(sl, h, x2+0.04, sy+0.03, col_w[j]-0.05, row_h-0.04, size=8, bold=True, color=WHITE)
for i, c in enumerate(customers[:10]):
    y = sy + row_h * (i + 1)
    m = sum(wt[c][:3])
    a = sum(wt[c][3:])
    tot = sum(wt[c])
    if m > 0:
        pct = int((a / m - 1) * 100)
        var = f"+{pct}%" if pct >= 0 else f"{pct}%"
        vcol = GREEN if pct >= 0 else RED
    else:
        var = "-"
        vcol = WHITE
    bg = ROW_ALT if i % 2 == 0 else CARD
    add_rect(sl, 0.3, y, tw, row_h, bg)
    vals = [str(i+1), c, str(m), str(a), str(tot), var]
    for j, (v, x2) in enumerate(zip(vals, col_x)):
        col = vcol if j == 5 else WHITE
        add_text(sl, v, x2+0.04, y+0.03, col_w[j]-0.05, row_h-0.04, size=8, color=col)

# ── SLIDE 4: HEATMAP 5 SEMANAS ───────────────────────────────────
sl = prs.slides.add_slide(blank)
set_bg(sl)
add_text(sl, 'Evolucao Semanal - Heatmap Top 20 x Periodos', 0.3, 0.1, 12, 0.42, size=16, bold=True)

fig, ax = plt.subplots(figsize=(11.5, 4.8))
matrix = np.array([wt[c] for c in customers], dtype=float)
max_val = matrix.max() if matrix.max() > 0 else 1
im = ax.imshow(matrix, cmap='Blues', aspect='auto', vmin=0, vmax=max_val)
ax.set_xticks(range(len(PERIODS)))
ax.set_xticklabels(PLABELS_SHORT, fontsize=10)
ax.set_yticks(range(len(customers)))
ax.set_yticklabels(customers, fontsize=8.5)
for i in range(len(customers)):
    for j in range(len(PERIODS)):
        v = int(matrix[i, j])
        if v > 0:
            txt_col = 'white' if matrix[i, j] > max_val * 0.55 else '#0f1117'
            ax.text(j, i, str(v), ha='center', va='center', fontsize=9, color=txt_col, fontweight='bold')
plt.colorbar(im, ax=ax, label='Qtd Outgoing', pad=0.01)
ax.set_facecolor('#1a1d27')
fig.patch.set_facecolor('#0f1117')
buf = fig_to_stream(fig)
sl.shapes.add_picture(buf, Inches(0.3), Inches(0.75), Inches(12.7), Inches(6.6))

# ── SLIDES 5+: DETALHE POR CLIENTE (top 10) ──────────────────────
for ci, c in enumerate(customers[:10]):
    sl = prs.slides.add_slide(blank)
    set_bg(sl)

    # header
    add_rect(sl, 0, 0, 13.33, 0.72, CARD)
    add_text(sl, f'#{ci+1}  Cliente: {c}', 0.3, 0.1, 9, 0.52, size=16, bold=True, color=ACCENT)
    add_text(sl, f'Total: {sum(wt[c]):,} outgoing  |  01/Mar - 16/Abr/2026',
             9.0, 0.15, 4.0, 0.4, size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    # KPI row
    for idx, (lbl, val) in enumerate(zip(PLABELS_SHORT, wt[c])):
        prev = wt[c][idx-1] if idx > 0 else None
        delta_str = None
        delta_up = None
        if prev and prev > 0:
            pct = int((val - prev) / prev * 100)
            arrow = 'up' if pct >= 0 else 'down'
            delta_str = f"{'▲' if pct>=0 else '▼'} {abs(pct)}% vs ant."
            delta_up = pct > 0
        kpi_box(sl, 0.3 + idx * 2.12, 0.82, 2.0, 0.95, lbl, f'{val:,}', delta_str, delta_up)

    # Line chart
    fig, ax = plt.subplots(figsize=(5.5, 2.8))
    ax.plot(PLABELS_SHORT, wt[c], color='#6366f1', linewidth=2.5, marker='o', markersize=6)
    ax.fill_between(PLABELS_SHORT, wt[c], alpha=0.15, color='#6366f1')
    for xi, yv in enumerate(wt[c]):
        if yv > 0:
            ax.annotate(str(yv), (xi, yv), textcoords='offset points', xytext=(0, 7),
                        ha='center', fontsize=8)
    ax.yaxis.grid(True)
    ax.set_axisbelow(True)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.set_title('Evolucao Total Semanal', fontsize=9, pad=6)
    buf = fig_to_stream(fig)
    sl.shapes.add_picture(buf, Inches(0.3), Inches(1.92), Inches(5.5), Inches(2.9))

    # Stacked by processo
    pw = proc_weekly(c)
    procs_list = list(pw.keys())
    fig2, ax2 = plt.subplots(figsize=(5.5, 2.8))
    bot = np.zeros(len(PERIODS))
    for pi, proc in enumerate(procs_list):
        vals = np.array(pw[proc])
        ax2.bar(PLABELS_SHORT, vals, bottom=bot,
                label=proc[:28], color=PALETTE_RGB[pi % len(PALETTE_RGB)], width=0.55)
        bot += vals
    ax2.yaxis.grid(True)
    ax2.set_axisbelow(True)
    ax2.spines['top'].set_visible(False)
    ax2.spines['right'].set_visible(False)
    ax2.set_title('Breakdown por Processo', fontsize=9, pad=6)
    handles, labels = ax2.get_legend_handles_labels()
    ax2.legend(handles[:5], [l[:22] for l in labels[:5]], fontsize=7, loc='upper right',
               facecolor='#1a1d27', edgecolor='#2a2d3a', labelcolor='#e2e8f0')
    buf2 = fig_to_stream(fig2)
    sl.shapes.add_picture(buf2, Inches(6.9), Inches(1.92), Inches(5.5), Inches(2.9))

    # Top motivos table
    agg = defaultdict(int)
    for period in PERIODS:
        for proc, cdus in data[c][period].items():
            for cdu_name, sols in cdus.items():
                for sol, q in sols.items():
                    agg[(proc, cdu_name, sol)] += q
    top_rows = sorted(agg.items(), key=lambda x: -x[1])[:7]

    add_text(sl, 'Top Motivos - Processo / CDU / Solucao', 0.3, 4.98, 12.7, 0.28, size=10, bold=True, color=MUTED)
    t_hdrs = ['Processo', 'CDU', 'Solucao', 'Qtd']
    t_cw   = [2.2, 2.0, 7.0, 0.75]
    t_cx   = [0.3]
    for w in t_cw[:-1]:
        t_cx.append(t_cx[-1] + w)
    rh = 0.21
    sy2 = 5.3
    add_rect(sl, 0.3, sy2, sum(t_cw), rh, ACCENT)
    for j, (h, x2) in enumerate(zip(t_hdrs, t_cx)):
        add_text(sl, h, x2+0.04, sy2+0.03, t_cw[j]-0.05, rh-0.04, size=8, bold=True, color=WHITE)
    for ri, ((proc, cdu_name, sol), q) in enumerate(top_rows):
        y = sy2 + rh * (ri + 1)
        if y > 7.28:
            break
        bg = ROW_ALT if ri % 2 == 0 else CARD
        add_rect(sl, 0.3, y, sum(t_cw), rh, bg)
        for j, (val, x2) in enumerate(zip([proc[:30], cdu_name[:28], sol[:80], str(q)], t_cx)):
            add_text(sl, val, x2+0.04, y+0.03, t_cw[j]-0.05, rh-0.04, size=8, color=WHITE)

# ── SAVE ─────────────────────────────────────────────────────────
out = "C:/claudinho/outgoing_dashboard.pptx"
prs.save(out)
print(f"Salvo: {out}  ({os.path.getsize(out)//1024} KB)  |  {len(prs.slides)} slides")
