import io, os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── colours ──────────────────────────────────────────────────────
BG      = RGBColor(0x0f, 0x11, 0x17)
CARD    = RGBColor(0x1a, 0x1d, 0x27)
ACCENT  = RGBColor(0x63, 0x66, 0xf1)
WHITE   = RGBColor(0xff, 0xff, 0xff)
MUTED   = RGBColor(0x94, 0xa3, 0xb8)
GREEN   = RGBColor(0x22, 0xc5, 0x5e)
RED     = RGBColor(0xef, 0x44, 0x44)
YELLOW  = RGBColor(0xf5, 0x9e, 0x0b)
ROW_ALT = RGBColor(0x1e, 0x21, 0x30)
ORANGE  = RGBColor(0xf9, 0x73, 0x16)

# ── dados coletados do BigQuery ───────────────────────────────────
DRIVER  = "Publicaciones Seller Dev"
S2_LBL  = "06–12/abr"
S1_LBL  = "13–19/abr"
REPORT_DATE = "20/04/2026"

NPS_S2_DRIVER = 65.54
NPS_S1_DRIVER = 63.54
VAR_DRIVER    = -2.00

s1 = {
    'Afiliados ML':               (220, 43, 284),
    'Gestión de Publicación':     (37,  10,  57),
    'PR - Técnica prohibida':     (39,   5,  54),
    'PR - Propiedad intelectual': (41,   5,  51),
    'Antes de publicar':          (44,   3,  50),
    'Potenciar Ventas':           (36,   6,  43),
    'PR - Artículos prohibidos':  (20,   7,  29),
    'Denuncia de usuario':        (10,   1,  12),
}
s2 = {
    'Afiliados ML':               (321, 39, 388),
    'Gestión de Publicación':     (62,  12,  80),
    'PR - Técnica prohibida':     (52,   5,  62),
    'PR - Propiedad intelectual': (44,  13,  59),
    'Antes de publicar':          (43,   7,  56),
    'Potenciar Ventas':           (33,  11,  52),
    'PR - Artículos prohibidos':  (31,   6,  39),
    'Denuncia de usuario':        (17,   3,  22),
}

def nps(p, d, s): return round(100.0 * (p - d) / s, 2) if s > 0 else 0.0

PROCESSOS_IMPACT = [
    ('Afiliados ML',               62.32, 72.68, -10.36, -5.38, 284),
    ('Gestión de Publicación',     47.37, 62.50, -15.13, -1.55,  57),
    ('PR - Técnica prohibida',     62.96, 75.81, -12.85, -0.92,  54),
    ('PR - Artículos prohibidos',  44.83, 64.10, -19.27, -0.98,  29),
    ('PR - Propiedad intelectual', 70.59, 52.54, +18.05, +1.25,  51),
    ('Antes de publicar',          82.00, 64.29, +17.71, +1.27,  50),
    ('Potenciar Ventas',           69.77, 42.31, +27.46, +1.73,  43),
    ('Denuncia de usuario',        75.00, 63.64, +11.36, +0.34,  12),
]

TEMAS_DETRACAO = [
    ('Comissão não creditada',  62, 'Vendas realizadas via link mas comissão cancelada/não\ncomputada. Regra de "auto-oferta" descobre-se só\nao perder a comissão.'),
    ('Atendimento sem solução', 28, 'Múltiplos atendentes sem capacidade de resolver.\nNenhum desbloqueio de conta, nenhuma reversão\nde comissão, nenhum encaminhamento eficaz.'),
    ('Regras confusas',         22, 'Link de 24h, restrição de venda para familiares/\nconhecidos, não pode usar WhatsApp individual.\nComunicação do programa insuficiente.'),
    ('Qualidade do atendimento', 8, 'Ligação caiu, silêncios longos, rep sem fone no mudo,\ntransferência prometida que não ocorreu.'),
]

INSIGHTS_TRANSCRICOES = [
    'Afiliados suspensos por "auto-oferta" descobrem a regra apenas ao serem punidos — não há alerta preventivo.',
    'REPs explicam as regras corretamente mas não têm autonomia para reverter bloqueios ou recrédito de comissão.',
    'Canal e-mail: cliente enviou informações 3x; cada REP pediu as mesmas informações novamente (sem acesso ao histórico).',
    'Ligações C2C com silêncio prolongado (até 345s) e quedas recorrentes — impacto direto na percepção de qualidade.',
    'Regra de auto-oferta (venda para pessoas vinculadas) é o principal trigger de insatisfação: afeta afiliados que divulgam para família/amigos.',
]

RECOMENDACOES = [
    ('Comunicação proativa das regras', 'CRÍTICO',
     'Enviar alertas preventivos a afiliados próximos do limiar de auto-oferta antes de\nbloqueio. Melhorar onboarding com regras de elegibilidade em linguagem simples.'),
    ('Fluxo dedicado Afiliados ML', 'ALTO',
     'Criar fila específica com REPs treinados e autonomia para reverter comissões\nindevidas e desbloquear contas em casos elegíveis, com SLA de 48h.'),
    ('Rever duração do link (24h)', 'MÉDIO',
     'Link de 24h inviabiliza estratégia de divulgação orgânica (posts, reels, stories).\nAvaliar extensão para 7–30 dias alinhado com a política de fraude.'),
    ('Histórico unificado no e-mail', 'MÉDIO',
     'REPs sem acesso ao histórico de contatos anteriores geram retrabalho e percepção\nde descaso. Integrar CRM ao canal e-mail para continuidade de atendimento.'),
]

# ── helpers ──────────────────────────────────────────────────────
plt.rcParams.update({
    'figure.facecolor': '#0f1117', 'axes.facecolor': '#1a1d27',
    'axes.edgecolor': '#2a2d3a',   'axes.labelcolor': '#94a3b8',
    'xtick.color': '#64748b',      'ytick.color': '#64748b',
    'text.color': '#e2e8f0',       'grid.color': '#1e2130',
    'grid.linewidth': 0.5,         'font.family': 'DejaVu Sans', 'font.size': 9,
})

def fig_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def set_bg(slide):
    bg = slide.background; fill = bg.fill; fill.solid()
    fill.fore_color.rgb = BG

def add_text(slide, text, l, t, w, h, size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color

def add_rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background(); return shape

def kpi_box(slide, l, t, w, h, label, value, delta=None, delta_bad=None):
    add_rect(slide, l, t, w, h, CARD)
    add_text(slide, label, l+0.08, t+0.06, w-0.16, 0.22, size=7, color=MUTED)
    add_text(slide, str(value), l+0.08, t+0.26, w-0.16, 0.38, size=16, bold=True, color=WHITE)
    if delta:
        dcol = RED if delta_bad else GREEN
        add_text(slide, delta, l+0.08, t+0.62, w-0.16, 0.22, size=8, color=dcol)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 0.08, 7.5, ACCENT)
add_rect(sl, 0.08, 0, 13.25, 0.5, CARD)
add_text(sl, 'DIAGNÓSTICO NPS', 0.5, 1.2, 12, 1.2, size=50, bold=True, color=ACCENT)
add_text(sl, 'Driver Publicaciones Seller Dev — Centro BR', 0.5, 2.65, 12, 0.65, size=22, color=WHITE)
add_text(sl, f'Comparativo Semanal: {S2_LBL}  →  {S1_LBL}', 0.5, 3.45, 12, 0.45, size=14, color=MUTED)
add_text(sl, 'Análise de variação por processo  |  Deep dive qualitativo: comentários + transcrições', 0.5, 4.0, 12, 0.35, size=11, color=MUTED)
add_rect(sl, 0.5, 4.55, 5.5, 0.06, RED)
add_text(sl, f'Gerado em {REPORT_DATE}', 0.5, 6.9, 6, 0.3, size=9, color=MUTED)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — RESULTADO DO DRIVER
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_text(sl, 'Publicaciones Seller Dev — Resultado Consolidado', 0.3, 0.1, 11, 0.42, size=15, bold=True)
add_text(sl, f'{S2_LBL}  →  {S1_LBL}  |  Centro BR', 0.3, 0.48, 8, 0.28, size=10, color=MUTED)

kpi_box(sl, 0.3, 0.8,  2.8, 1.0, f'NPS S2 ({S2_LBL})', f'{NPS_S2_DRIVER:.2f}%')
kpi_box(sl, 3.3, 0.8,  2.8, 1.0, f'NPS S1 ({S1_LBL})', f'{NPS_S1_DRIVER:.2f}%')
kpi_box(sl, 6.3, 0.8,  2.8, 1.0, 'Variação (pp)',      f'{VAR_DRIVER:+.2f} pp', '▼ queda na semana', True)
kpi_box(sl, 9.3, 0.8,  3.7, 1.0, 'Volume S1',          '587 surveys',           '↓ 182 vs S2 (769)', True)

# waterfall chart por processo
fig, ax = plt.subplots(figsize=(12.0, 3.8))
procs_w = ['Afiliados\nML', 'Gestión de\nPublicación', 'PR-Técnica\nproh.', 'PR-Art.\nproh.',
           'PR-Prop.\nintl.', 'Antes de\npublicar', 'Potenciar\nVentas', 'Denuncia\nusuario']
impacts_w = [-5.38, -1.55, -0.92, -0.98, 1.25, 1.27, 1.73, 0.34]
cols_w = ['#ef4444' if v < 0 else '#22c55e' for v in impacts_w]
bars = ax.bar(procs_w, impacts_w, color=cols_w, width=0.6, edgecolor='#2a2d3a', linewidth=0.5)
ax.axhline(0, color='#64748b', linewidth=0.8, linestyle='--')
for bar, val in zip(bars, impacts_w):
    ax.text(bar.get_x()+bar.get_width()/2, val + (0.12 if val>0 else -0.28),
            f'{val:+.2f}', ha='center', va='bottom' if val>0 else 'top',
            fontsize=9, fontweight='bold', color='white')
ax.set_ylabel('Impacto no Driver (pp)', fontsize=9)
ax.set_title(f'Impacto MIX+NETO por Processo  |  {S2_LBL} → {S1_LBL}', fontsize=10, pad=6)
ax.yaxis.grid(True, alpha=0.4); ax.set_axisbelow(True)
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
sl.shapes.add_picture(fig_stream(fig), Inches(0.3), Inches(2.05), Inches(12.7), Inches(4.1))

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — TABELA DE PROCESSOS
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_text(sl, 'Variação por Processo — Publicaciones Seller Dev', 0.3, 0.1, 11, 0.42, size=15, bold=True)
add_text(sl, f'{S2_LBL}  →  {S1_LBL}  |  Centro BR  |  Ordenado por impacto', 0.3, 0.48, 10, 0.28, size=10, color=MUTED)

hdrs  = ['Processo', 'NPS S2', 'NPS S1', 'VAR (pp)', 'Surveys S1', 'Impacto Driver']
col_w = [3.4, 1.25, 1.25, 1.25, 1.35, 1.65]
col_x = [0.3]
for w in col_w[:-1]: col_x.append(col_x[-1] + w)
row_h = 0.38
sy = 0.82

add_rect(sl, 0.3, sy, sum(col_w), row_h, ACCENT)
for j, (h, x) in enumerate(zip(hdrs, col_x)):
    add_text(sl, h, x+0.06, sy+0.08, col_w[j]-0.1, row_h-0.1, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

rows_data = sorted(PROCESSOS_IMPACT, key=lambda r: r[4])
for i, (proc, n1, n2, var, imp, vol) in enumerate(rows_data):
    y = sy + row_h*(i+1)
    bg = ROW_ALT if i % 2 == 0 else CARD
    add_rect(sl, 0.3, y, sum(col_w), row_h, bg)
    vals  = [proc, f'{n2:.2f}%', f'{n1:.2f}%', f'{var:+.2f}', str(vol), f'{imp:+.2f} pp']
    colors = [WHITE, WHITE, WHITE,
              RED if var < 0 else GREEN,
              WHITE,
              RED if imp < 0 else GREEN]
    for j, (v, x, c) in enumerate(zip(vals, col_x, colors)):
        add_text(sl, v, x+0.06, y+0.09, col_w[j]-0.1, row_h-0.12, size=9, color=c,
                 align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

note_y = sy + row_h*(len(rows_data)+1) + 0.12
add_text(sl, '* Impacto = contribuição de cada processo na variação total do driver (MIX + NETO). Processos com volume < 5 surveys excluídos.',
         0.3, note_y, 12.7, 0.28, size=7, color=MUTED)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — FOCO AFILIADOS ML
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_rect(sl, 0, 0, 0.08, 7.5, RED)
add_text(sl, 'Deep Dive — Afiliados ML', 0.3, 0.1, 11, 0.42, size=15, bold=True, color=WHITE)
add_text(sl, 'Processo de maior impacto negativo no driver', 0.3, 0.48, 10, 0.28, size=10, color=RED)

kpi_box(sl, 0.3, 0.82, 2.2, 1.0, f'NPS S2 ({S2_LBL})',  '72,68%')
kpi_box(sl, 2.7, 0.82, 2.2, 1.0, f'NPS S1 ({S1_LBL})',  '62,32%', '▼ -10,36 pp', True)
kpi_box(sl, 5.1, 0.82, 2.2, 1.0, 'Impacto no Driver',  '-5,38 pp', '80% da queda total', True)
kpi_box(sl, 7.5, 0.82, 2.2, 1.0, 'Detratores S1',      '43 de 284', f'15,1% taxa det.', True)
kpi_box(sl, 9.9, 0.82, 3.1, 1.0, 'Volume S1 vs S2',    '284 → 388', '↓ 104 surveys (-26,8%)', True)

# NPS S2 vs S1 chart
fig, axes = plt.subplots(1, 2, figsize=(11.5, 3.2))
ax1, ax2 = axes

# Bar chart NPS comparison
cats = [f'S2\n{S2_LBL}', f'S1\n{S1_LBL}']
vals_nps = [72.68, 62.32]
cols = ['#6366f1', '#ef4444']
bars1 = ax1.bar(cats, vals_nps, color=cols, width=0.5, edgecolor='#2a2d3a')
ax1.axhline(65.54, color='#64748b', linewidth=1, linestyle='--', label='NPS Total Driver S1')
for bar, val in zip(bars1, vals_nps):
    ax1.text(bar.get_x()+bar.get_width()/2, val+0.8, f'{val:.2f}%',
             ha='center', fontsize=11, fontweight='bold', color='white')
ax1.set_ylim(55, 80); ax1.yaxis.grid(True, alpha=0.4); ax1.set_axisbelow(True)
ax1.set_title('NPS Afiliados ML — S2 vs S1', fontsize=10, pad=5)
ax1.spines['top'].set_visible(False); ax1.spines['right'].set_visible(False)

# Donut detractors/promoters S1
sizes = [220, 43, 21]
colors_d = ['#22c55e', '#ef4444', '#94a3b8']
labels_d = ['Promotores\n220', 'Detratores\n43', 'Neutros\n21']
wedges, texts = ax2.pie(sizes, colors=colors_d, startangle=90, wedgeprops=dict(width=0.55))
ax2.text(0, 0, f'NPS\n62,32%', ha='center', va='center', fontsize=11, fontweight='bold', color='white')
ax2.legend(wedges, labels_d, loc='lower center', bbox_to_anchor=(0.5, -0.18),
           ncol=3, fontsize=8, framealpha=0, labelcolor='white')
ax2.set_title('Distribuição Surveys S1 (13–19/abr)', fontsize=10, pad=5)

sl.shapes.add_picture(fig_stream(fig), Inches(0.3), Inches(2.0), Inches(12.7), Inches(3.5))

add_text(sl, '⚡ DIAGNÓSTICO: Afiliados ML concentra 48% do volume do driver e registrou queda de 10,36 pp. Principal driver da retração semanal de -2,00 pp.',
         0.3, 5.7, 12.7, 0.45, size=10, bold=True, color=YELLOW)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — TEMAS DE DETRAÇÃO
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_text(sl, 'Temas de Detração — Afiliados ML', 0.3, 0.1, 11, 0.42, size=15, bold=True)
add_text(sl, f'Análise de comentários livres de detratores (NPS 0–6)  |  Semana {S1_LBL}  |  amostra 85 comentários', 0.3, 0.48, 12, 0.28, size=9, color=MUTED)

# Horizontal bar chart
fig, ax = plt.subplots(figsize=(5.8, 2.8))
temas_lbl = [t[0] for t in TEMAS_DETRACAO]
temas_pct = [t[1] for t in TEMAS_DETRACAO]
t_cols = ['#ef4444', '#f97316', '#f59e0b', '#94a3b8']
bars = ax.barh(temas_lbl[::-1], temas_pct[::-1], color=t_cols[::-1], height=0.55, edgecolor='#2a2d3a')
for bar, val in zip(bars, temas_pct[::-1]):
    ax.text(val+0.5, bar.get_y()+bar.get_height()/2, f'{val}%',
            va='center', fontsize=9, fontweight='bold', color='white')
ax.set_xlim(0, 80); ax.xaxis.grid(True, alpha=0.3); ax.set_axisbelow(True)
ax.set_xlabel('% citações nos comentários')
ax.spines['top'].set_visible(False); ax.spines['right'].set_visible(False)
ax.set_title('Temas mencionados por detratores (%)', fontsize=10, pad=5)
sl.shapes.add_picture(fig_stream(fig), Inches(0.3), Inches(0.9), Inches(6.0), Inches(3.1))

# Detail cards
card_y = 0.85
for i, (tema, pct, descr) in enumerate(TEMAS_DETRACAO):
    cy = card_y + i * 1.52
    col = [RED, ORANGE, YELLOW, MUTED][i]
    add_rect(sl, 6.5, cy, 6.5, 1.38, CARD)
    add_rect(sl, 6.5, cy, 0.06, 1.38, col)
    add_text(sl, f'{tema} — {pct}% das citações', 6.7, cy+0.1, 6.2, 0.28, size=9, bold=True, color=WHITE)
    add_text(sl, descr, 6.7, cy+0.38, 6.1, 0.95, size=8, color=MUTED)

add_text(sl, '* Percentuais baseados em frequência de menção nos comentários livres. Um comentário pode citar múltiplos temas.',
         0.3, 7.15, 12.7, 0.25, size=7, color=MUTED)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — DIAGNÓSTICO DAS TRANSCRIÇÕES
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_text(sl, 'Análise de Transcrições — Padrões de Atendimento', 0.3, 0.1, 11, 0.42, size=15, bold=True)
add_text(sl, 'Cruzamento comentários × transcrições USER-REP  |  Afiliados ML  |  Semana 13–19/abr', 0.3, 0.48, 12, 0.28, size=9, color=MUTED)

add_text(sl, '🔍 Achados nas Transcrições', 0.3, 0.78, 12, 0.32, size=11, bold=True, color=ACCENT)
icons = ['🚫', '⚙️', '📧', '📵', '⚠️']
iy = 1.15
for idx, (icon, insight) in enumerate(zip(icons, INSIGHTS_TRANSCRICOES)):
    add_rect(sl, 0.3, iy + idx*1.08, 12.5, 1.0, CARD)
    add_text(sl, icon, 0.45, iy+idx*1.08+0.28, 0.45, 0.5, size=18)
    add_text(sl, insight, 1.0, iy+idx*1.08+0.1, 11.6, 0.82, size=9, color=WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — CAUSA RAIZ
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_text(sl, 'Diagnóstico — Causa Raiz da Queda', 0.3, 0.1, 11, 0.42, size=15, bold=True)
add_text(sl, 'Afiliados ML  |  -10,36 pp  |  Semana 13–19/abr', 0.3, 0.48, 10, 0.28, size=10, color=RED)

# 3 causa boxes
causas = [
    ('Sistêmico', 'Regras de comissionamento\n(auto-oferta, link 24h)',
     'Afiliados desconhecem restrições até perder a comissão.\nOnboarding não esclarece casos de inelegibilidade.\nDescoberta tardia = frustração e percepção de enganação.'),
    ('Operacional', 'Atendimento sem autonomia\nresolutiva',
     'REPs explicam as regras mas não podem reverter bloqueios\nnem recrédito de comissões. Cada contato gera nova frustração\nsem avançar na resolução do problema real.'),
    ('Experiência', 'Produto percebido como\n"propaganda enganosa"',
     'Afiliados que trabalham para divulgar o ML e não recebem\ncomissão migram para outras plataformas e geram NPS 0.\nPerda de canal orgânico de aquisição de clientes.'),
]
caus_cols = [RED, ORANGE, YELLOW]
for i, (titulo, subtit, descr) in enumerate(causas):
    cx = 0.3 + i * 4.35
    add_rect(sl, cx, 0.82, 4.1, 3.8, CARD)
    add_rect(sl, cx, 0.82, 4.1, 0.06, caus_cols[i])
    add_text(sl, titulo, cx+0.15, 0.95, 3.8, 0.32, size=12, bold=True, color=caus_cols[i])
    add_text(sl, subtit, cx+0.15, 1.28, 3.8, 0.55, size=10, bold=True, color=WHITE)
    add_text(sl, descr,  cx+0.15, 1.88, 3.8, 2.65, size=8.5, color=MUTED)

add_text(sl, '📌 CONCLUSÃO', 0.3, 4.78, 12.7, 0.32, size=11, bold=True, color=ACCENT)
add_rect(sl, 0.3, 5.15, 12.7, 1.82, CARD)
add_text(sl,
    'A queda de -10,36 pp em Afiliados ML é estrutural, não pontual. O funil de problemas começa antes do atendimento:\n'
    'afiliados chegam ao CX já insatisfeitos pois o produto (programa de afiliados) não entregou o prometido. O atendimento\n'
    'amplifica a insatisfação ao não ter capacidade de resolver. Ação imediata deve focar em comunicação preventiva das regras\n'
    'e criação de fluxo com autonomia resolutiva para REPs que atendem este processo.',
    0.5, 5.25, 12.3, 1.65, size=9.5, color=WHITE)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — RECOMENDAÇÕES
# ════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank); set_bg(sl)
add_rect(sl, 0, 0, 13.33, 0.62, CARD)
add_text(sl, 'Recomendações Acionáveis', 0.3, 0.1, 11, 0.42, size=15, bold=True)
add_text(sl, 'Afiliados ML — Publicaciones Seller Dev  |  Priorizadas por criticidade', 0.3, 0.48, 12, 0.28, size=9, color=MUTED)

prio_col = {'CRÍTICO': RED, 'ALTO': ORANGE, 'MÉDIO': YELLOW}
layout = [(0.3, 0.82, 6.1, 3.1), (6.6, 0.82, 6.4, 3.1),
          (0.3, 4.05, 6.1, 3.1), (6.6, 4.05, 6.4, 3.1)]
for i, (titulo, prio, descr) in enumerate(RECOMENDACOES):
    x, y, w, h = layout[i]
    col = prio_col.get(prio, MUTED)
    add_rect(sl, x, y, w, h, CARD)
    add_rect(sl, x, y, w, 0.06, col)
    prio_tag = f'  {prio}'
    add_text(sl, prio_tag, x+0.12, y+0.1, 1.5, 0.24, size=8, bold=True, color=col)
    add_text(sl, titulo, x+0.12, y+0.35, w-0.2, 0.32, size=10, bold=True, color=WHITE)
    add_text(sl, descr,  x+0.12, y+0.72, w-0.2, 2.25, size=8.5, color=MUTED)

# ════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════
out = "C:/claudinho/nps_publicaciones_s1_2026.pptx"
prs.save(out)
print(f"Salvo: {out}  ({os.path.getsize(out)//1024} KB)  |  {len(prs.slides)} slides")
